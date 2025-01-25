import pandas as pd
from pathlib import Path
import random
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import MSO_AUTO_SIZE
import matplotlib.pyplot as plt
from datetime import datetime
import json

def clean_data(input_dir=Path("mock_tables"), output_dir=Path("clean_tables")):
    """
    Clean supplier data by filtering specific categories from input Excel files.
    
    Args:
        input_dir (Path, optional): Directory containing input files. Defaults to 'mock_tables'.
        output_dir (Path, optional): Directory for output files. Defaults to 'clean_tables'.
    """


    # Check if required files exist
    required_files = [
        "mock_Supplier_fact_sheet_IT_spend_2024.xlsx",
        "mock_Supplier_fact_sheet_Contracting_report.xlsx",
        "mock_Supplier_fact_sheet_sourcing_event_participation.xlsx"
    ]
    
    missing_files = [file for file in required_files if not (input_dir / file).exists()]
    
    if missing_files:
        raise FileNotFoundError(
            f"Missing required files: {', '.join(missing_files)}\n"
            "Please run mock_tables/mock_tables.py first to generate the required files."
        )

    # If no arguments provided, input_dir and output_dir will use default Path values
    # Create output dir if it doesn't exist
    output_dir.mkdir(exist_ok= True)

    # Load the data
    it_spend = pd.read_excel(input_dir / "mock_Supplier_fact_sheet_IT_spend_2024.xlsx")
    contracting_report = pd.read_excel(input_dir / "mock_Supplier_fact_sheet_Contracting_report.xlsx")
    sourcing_event = pd.read_excel(input_dir / "mock_Supplier_fact_sheet_sourcing_event_participation.xlsx")

    # Filter data based on the required categories
    filtered_it_spend = it_spend[it_spend["Category"].isin(["IT Hardware", "Telecoms and Network"])]
    filtered_contracting_report = contracting_report[
        contracting_report["[PCW] OneProcurement Category"] == "IT Infrastructure"
    ]
    filtered_sourcing_event = sourcing_event[
        sourcing_event["[SPRJ] OneProcurement Category"] == "IT Infrastructure"
    ]

    # Save the cleaned data
    filtered_it_spend.to_excel(output_dir / "cleaned_Supplier_fact_sheet_IT_spend_2024.xlsx", index=False)
    filtered_contracting_report.to_excel(output_dir / "cleaned_Supplier_fact_sheet_Contracting_report.xlsx", index=False)
    filtered_sourcing_event.to_excel(output_dir / "cleaned_Supplier_fact_sheet_sourcing_event_participation.xlsx", index=False)
    
    return filtered_it_spend, filtered_contracting_report, filtered_sourcing_event

def extract_unique_vendors(it_spend, contracting_report, sourcing_event):
    """
    Extract a list of unique vendors across the three datasets.
    
    Args:
        it_spend (DataFrame): Cleaned IT spend data.
        contracting_report (DataFrame): Cleaned contracting report data.
        sourcing_event (DataFrame): Cleaned sourcing event data.

    Returns:
        set: A set of unique vendor names.
    """
    # Extract unique vendors from IT spend data
    vendors_it_spend = set(it_spend["Vendor Name"].unique())
    
    # Extract unique vendors from Contracting Report
    vendors_contracting = set(contracting_report["[PCW]Affected Parties (Supplier Name (L1))"].unique())
    
    # Extract unique vendors from Sourcing Event data
    vendors_sourcing = set(sourcing_event["[SPT]Supplier (Supplier Name (L1))"].unique())
    
    # Combine all unique vendors into a single set to remove duplicates
    unique_vendors = vendors_it_spend.union(vendors_contracting, vendors_sourcing)
    
    return unique_vendors

def simulate_historical_spend(it_spend):
    """
    Simulate spend data for 2022 and 2023 based on 2024 spend.
    
    Args:
        it_spend (DataFrame): Cleaned IT spend data with 2024 spend.

    Returns:
        DataFrame: Updated DataFrame with simulated 2022, 2023 spend.
    """
    # Add 2022 and 2023 spend columns
    simulated_data = it_spend.copy()
    simulated_data["Spend 2023 (€m)"] = simulated_data["Spend 2024 (€m)"] * (1 + simulated_data.apply(lambda x: random.uniform(-0.2, 0.2), axis=1))
    simulated_data["Spend 2022 (€m)"] = simulated_data["Spend 2023 (€m)"] * (1 + simulated_data.apply(lambda x: random.uniform(-0.2, 0.2), axis=1))
    return simulated_data



def generate_key_contracts_table(slide, contracting_data, position, row_height_cm=0.6):
    """
    Add a Key Contracts table to the slide.
    
    Args:
        slide: A PowerPoint slide in a presentation
        contracting_data (DataFrame): Filtered contracting data for the vendor.
        position (tuple): Coordinates (left, top) for the table.
    """

    # Define table headers
    headers = ["ID", "Name", "Short description", "Term", "Exp. date", "TCV (€m)"]
    cols = len(headers)
    
    # Define column widths (in cm) - total should sum to 19
    col_widths = [2.5, 3.5, 5, 3, 2.5, 2.5]  # Adjusted to sum to 19 cm
    
    # Determine the number of rows
    rows = min(len(contracting_data), 4) if not contracting_data.empty else 1

    # Create the table with specific dimensions
    table = slide.shapes.add_table(
        rows + 1,  # +1 for header row
        cols,
        Cm(position[0]),
        Cm(position[1]),
        Cm(19),    # Total width changed to 19 cm
        Cm(rows * row_height_cm)  # Total height based on number of rows
    ).table

    # Set fixed row height for all rows
    for row_idx in range(rows + 1):
        table.rows[row_idx].height = Cm(row_height_cm)

    # Set column widths
    for i, width in enumerate(col_widths):
        table.columns[i].width = Cm(width)

    # Format headers
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        # Set font for header
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(9)
        paragraph.font.name = "Arial"
        paragraph.font.bold = True

    # Reset the index of the contracting data
    contracting_data = contracting_data.reset_index(drop=True)

    # Populate the table rows
    if not contracting_data.empty:
        for row_idx in range(min(rows, len(contracting_data))):
            row = contracting_data.iloc[row_idx]
            # Populate and format each cell
            cells = [
                (str(row.get("[PCW] Contract Id", ""))[:8], 0),
                (str(row.get("[PCW]Contract (Contract)", ""))[:15], 1),
                (str(row.get("[PCW] Description", ""))[:25], 2),
                (f"{row.get('[PCW]Contract (Effective Date)', '').year}-{row.get('[PCW]Contract (Expiration Date)', '').year}" 
                 if pd.notna(row.get('[PCW]Contract (Effective Date)', '')) else "N/A", 3),
                (str(row.get('[PCW]Contract (Expiration Date)', '').date()) 
                 if pd.notna(row.get('[PCW]Contract (Expiration Date)', '')) else "N/A", 4),
                (str(round(row.get("sum(Contract Amount) (€m)", 0), 1)), 5)
            ]
            
            for text, col in cells:
                cell = table.cell(row_idx + 1, col)
                cell.text = text
                # Set font for cell
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(9)
                paragraph.font.name = "Arial"
    else:
        # Fill empty row with formatted cells
        for col_idx in range(cols):
            cell = table.cell(1, col_idx)
            cell.text = ""
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(9)
            paragraph.font.name = "Arial"

def generate_planned_projects_table(slide, sourcing_data, position, row_height_cm=0.6):
#     """
#     Add a Planned Projects table to the slide.
    
#     Args:
#         slide: The PowerPoint slide.
#         sourcing_data (DataFrame): Filtered sourcing event data for the vendor.
#         position (tuple): Coordinates (left, top) for the table.
#     """
    # Define headers and column widths
    headers = ["Project", "Name", "Short Description", "Value (€m)"]
    col_widths = [4, 5, 7, 3]  # Total 19 cm
    cols = len(headers)
    
    # Limit to 3 rows maximum
    rows = min(len(sourcing_data), 3) if not sourcing_data.empty else 1

    # Create table with specific dimensions
    table = slide.shapes.add_table(
        rows + 1,
        cols,
        Cm(position[0]),
        Cm(position[1]),
        Cm(19),    # Total width changed to 19 cm
        Cm(rows * row_height_cm)
    ).table

    # Set fixed row height for all rows
    for row_idx in range(rows + 1):
        table.rows[row_idx].height = Cm(row_height_cm)

    # Set column widths
    for i, width in enumerate(col_widths):
        table.columns[i].width = Cm(width)

    # Format headers
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        # Set font for header
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(9)
        paragraph.font.name = "Arial"
        paragraph.font.bold = True

    # Reset the index of the sourcing data
    sourcing_data = sourcing_data.reset_index(drop=True)

    # Populate table rows
    if not sourcing_data.empty:
        for row_idx in range(min(rows, len(sourcing_data))):
            row = sourcing_data.iloc[row_idx]
            cells = [
                (str(row["[SPRJ]Project (Project Id)"])[:8], 0),
                (str(row["[SPRJ]Project (Project Name)"])[:20], 1),
                (str(row["Short Description"])[:30], 2),
                (str(round(row["sum(Baseline Spend) (€m)"], 1)), 3)
                

            ]
            
            for text, col in cells:
                cell = table.cell(row_idx + 1, col)
                cell.text = text
                # Set font for cell
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(9)
                paragraph.font.name = "Arial"
    else:
        # Fill empty row with formatted cells
        for col_idx in range(cols):
            cell = table.cell(1, col_idx)
            cell.text = ""
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(9)
            paragraph.font.name = "Arial"


# Generate a bar chart for IT spend overview.
def generate_it_spend_chart(vendor_name, spend_data, output_dir="plots"):
    """
    Generate a bar chart for IT spend overview.
    
    Args:
        vendor_name (str): Vendor name.
        spend_data (Series): Historical spend data (2022, 2023, 2024).
    """
    # Set up plots directory
    output_dir = Path(output_dir)
    output_dir.mkdir(exist_ok=True)

    # Define coordinated blue color palette
    blue_colors = ['#1f77b4',  # deep blue
                  '#6baed6',  # medium blue
                  '#bdd7e7']  # light blue
    
    years = ["2022", "2023", "2024"]
    spends = spend_data[["Spend 2022 (€m)", "Spend 2023 (€m)", "Spend 2024 (€m)"]]
    
    plt.bar(years, spends, color=blue_colors)
    plt.title(f"IT Spend Overview (€m incl. VAT) - {vendor_name}")
    plt.xlabel("Year")
    plt.ylabel("Spend (€m)")

    # Save to plots directory
    output_path = output_dir / f"{vendor_name}_spend_chart.png"
    plt.savefig(output_path)
    plt.close()
    

def add_image_to_slide(slide, image_path, position, size):
    """
    Add an image (e.g., chart) to the specified position on a slide.
    
    Args:
        slide: The PowerPoint slide object.
        image_path (str): Path to the image file.
        position (tuple): Position as (left, top) in centimeters.
        size (tuple): Size as (width, height) in centimeters.
    """
    left = Cm(position[0])
    top = Cm(position[1])
    width = Cm(size[0])
    height = Cm(size[1])
    slide.shapes.add_picture(image_path, left, top, width, height)

# # Search and replace text in PowerPoint while preserving formatting
# Issue: cannot hanlde multi-line replacements 
# def search_and_replace(search_str, repl_str, prs):
#     """"search and replace text in PowerPoint while preserving formatting"""

#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if shape.has_text_frame:
#                 if(shape.text.find(search_str))!=-1:
#                     text_frame = shape.text_frame
#                     cur_text = text_frame.paragraphs[0].runs[0].text
#                     new_text = cur_text.replace(str(search_str), str(repl_str))
#                     text_frame.paragraphs[0].runs[0].text = new_text
 
def search_and_replace(search_str, repl_str, prs):
    """
    Search and replace text in PowerPoint while preserving formatting.
    Handles multi-line replacements while maintaining original font formatting.
    """
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                if search_str in text_frame.text:
                    if len(text_frame.paragraphs) > 0:
                        first_para = text_frame.paragraphs[0]
                        if search_str in first_para.text:
                            # Store original formatting from the first run
                            original_font = first_para.runs[0].font
                            
                            # Clear existing paragraphs
                            while len(text_frame.paragraphs) > 1:
                                tr = text_frame._element.get_or_add_trPr()
                                tr.remove(text_frame.paragraphs[-1]._element)
                            
                            # Split replacement text into lines
                            lines = repl_str.split('\n')
                            
                            # Replace text in first paragraph while preserving format
                            first_para.text = first_para.text.replace(search_str, lines[0])
                            for run in first_para.runs:
                                run.font.name = original_font.name
                                run.font.size = original_font.size
                                run.font.bold = original_font.bold
                                run.font.italic = original_font.italic
                            
                            # Add new paragraphs for remaining lines
                            for line in lines[1:]:
                                p = text_frame.add_paragraph()
                                run = p.add_run()
                                run.text = line
                                # Copy all font attributes from original
                                run.font.name = original_font.name
                                run.font.size = original_font.size
                                run.font.bold = original_font.bold
                                run.font.italic = original_font.italic


def replace_vendor_placeholders(prs, vendor_name, vendor_data):
    """
    Replace all vendor-related placeholders in the presentation with actual data.
    
    Args:
        prs: PowerPoint presentation object
        vendor_name (str): Name of the vendor
        vendor_data (dict): Vendor data from all_vendors_data for a specific vendor
    """
    # Generate timestamp
    timestamp = f"VENDOR FACT SHEET - AS AT {datetime.now().strftime('%d.%m.%Y')}"

    # Format financials with bullet points
    financials_text = (
        f"• Revenue: {vendor_data['Financials']['Revenue']}\n"
        f"• Market Cap: {vendor_data['Financials']['MarketCap']}\n"
        f"• Growth Rate: {vendor_data['Financials']['GrowthRate']}"
    )
    key_account_managers = "\n".join(f"• {manager}" for manager in vendor_data["KeyAccountManagers"])
    key_stakeholders = "\n".join(f"• {stakeholder}" for stakeholder in vendor_data["KeyStakeholders"])
    
    # Define placeholder mappings
    replacements = {
        "[Timestamp]": timestamp,
        "[Vendor Name]": vendor_name,
        "[KeyAccountManagers]":key_account_managers,
        "[KeyStakeholders]": key_stakeholders,
        "[Financials]": financials_text,
        "[MarketTrends]": vendor_data["MarketTrends"],
        "[Strategy]": vendor_data["Strategy"],
        "[Msg]": vendor_data["Msg"]
    }
    
    # Replace each placeholder
    for placeholder, value in replacements.items():
        search_and_replace(placeholder, value, prs)


def generate_vendor_fact_sheets(template_path, output_dir, vendors, it_spend, contracting_report, sourcing_event):
    """
    Generate vendor fact sheets based on a template.
    
    Args:
        template_path (str): Path to the PowerPoint template.
        output_path (str): Path to save the generated PowerPoint file.
        vendors (set): Unique vendor names.
        it_spend (DataFrame): IT spend data.
        contracting_report (DataFrame): Contracting report data.
        sourcing_event (DataFrame): Sourcing event data.
    """
     # Load vendor data from JSON file
    with open('all_vendors_data.json', 'r') as f:
        all_vendors_data = json.load(f)
    
    # Debug 1 
    print(f"Total vendors: {len(vendors)}")
    print(f"IT spend shape: {it_spend.shape}")
    print(f"Sample vendors from it_spend: {it_spend['Vendor Name'].unique()[:5]}")
    
    # Create output directory at the start of the function
    output_dir = Path(output_dir)
    output_dir.mkdir(exist_ok=True)

    # Load the PowerPoint template
    prs = Presentation(template_path)
    vendor_slide = prs.slides[0]  # Access the pre-existing slide in the template
  
    for vendor in vendors:
        # Debug 2: print the current vendor
        print(f"\nProcessing vendor: {vendor}")
       
        # For each vendor, create a new presentation from the template
        prs = Presentation(template_path)
        vendor_slide = prs.slides[0]  # Get the first slide from the template
        
        # Filter data for the current vendor
        vendor_spend = it_spend[it_spend["Vendor Name"] == vendor]
        vendor_contracts = contracting_report[contracting_report["[PCW]Affected Parties (Supplier Name (L1))"] == vendor]
        vendor_projects = sourcing_event[sourcing_event["[SPT]Supplier (Supplier Name (L1))"] == vendor]

        # Debug information
        print(f"Vendor spend data shape: {vendor_spend.shape}")
        print(f"Vendor contracts data shape: {vendor_contracts.shape}")
        print(f"Vendor projects data shape: {vendor_projects.shape}")
        
        # # Replace the placeholder text in the template
        # timestamp = f"VENDOR FACT SHEET - AS AT {datetime.now().strftime('%d.%m.%Y')}"
        # search_and_replace("[Timestamp]", timestamp, prs)
        # search_and_replace("[Vendor Name]", vendor, prs)

        # Add IT Spend Overview Chart
        generate_it_spend_chart(vendor, vendor_spend.iloc[0], output_dir="plots")
        
        #left margin is 12cm, top margin is 5cm, width is 7cm, height is 5cm
        add_image_to_slide(vendor_slide, f"plots/{vendor}_spend_chart.png", position=(13, 4.0), size=(9.2, 5.5))
        
        # Add Key Contracts Table
        generate_key_contracts_table(vendor_slide, vendor_contracts, (13.5, 10.6))
        
        # Add Planned Projects Table
        generate_planned_projects_table(vendor_slide, vendor_projects, (13.5, 15.3))
    
        # Replace vendor-specific placeholders
        if vendor in all_vendors_data:
            replace_vendor_placeholders(prs, vendor, all_vendors_data[vendor])  
        
         # Save individual presentation for this vendor
        output_path = output_dir / f"{vendor}_Vendor_Fact_Sheet.pptx"
        prs.save(output_path)
    
    print(f"Vendor fact sheets saved to {str(output_dir)}")


def main():
    # Clean the data
    it_spend, contracting_report, sourcing_event = clean_data()
    print("Data cleaning completed!")

    # Initialize the simulated historical spend data
    simulated_spend = simulate_historical_spend(it_spend)

    # Extract the unique vendors
    unique_vendors = extract_unique_vendors(it_spend, contracting_report, sourcing_event)
    print(f"Unique vendors extracted: {len(unique_vendors)}")
    print(unique_vendors)

    # Generate vendor fact sheets
    generate_vendor_fact_sheets(
        template_path="pptx_template/vendor_template.pptx",
        output_dir="output",
        vendors=unique_vendors,
        it_spend=simulated_spend,
        contracting_report=contracting_report,
        sourcing_event=sourcing_event
    )

if __name__ == "__main__":
    main()